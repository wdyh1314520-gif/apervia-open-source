# Split from app3_parts/tools/file_registry_edit_tools_part.py.
# Purpose: document readers and OCR helpers.
# Loaded by app3.py via _exec_split_file(...), sharing the original global namespace.

def read_pptx(raw: bytes, max_slides: int = 30, max_shapes_per_slide: int = 200) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise RuntimeError("需要 python-pptx（pip install python-pptx）") from e

    prs = Presentation(io.BytesIO(raw))
    out = []
    for idx, slide in enumerate(prs.slides[:max_slides], start=1):
        out.append(f"# Slide {idx}")
        count = 0
        for shape in slide.shapes:
            if count >= max_shapes_per_slide:
                out.append("[已截断：对象过多]")
                break
            count += 1
            txt = ""
            try:
                if hasattr(shape, "text") and shape.text:
                    txt = str(shape.text).strip()
            except Exception:
                txt = ""
            if txt:
                out.append(txt)
                continue
            # table text
            try:
                table = getattr(shape, "table", None)
                if table is not None:
                    for row in table.rows:
                        vals = [str(cell.text or "").strip() for cell in row.cells]
                        if any(vals):
                            out.append("\t".join(vals).rstrip())
            except Exception:
                pass
        out.append("")
    return "\n".join(out).strip()


def read_archive_bundle(raw: bytes, ext: str, max_entries: int = 20, max_each_bytes: int = 800000, max_total_chars: int = 24000) -> str:
    import zipfile
    ext = (ext or "").lower()
    if ext != ".zip":
        raise RuntimeError("当前仅支持自动解析 .zip，其他压缩格式保留下载。")

    text_like_ext = {
        ".txt", ".md", ".json", ".jsonl", ".csv", ".tsv", ".py", ".c", ".cc", ".cpp", ".cxx", ".h", ".hpp",
        ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".mts", ".cts", ".java", ".go", ".rs", ".php", ".rb", ".swift", ".kt", ".cs",
        ".sql", ".yaml", ".yml", ".xml", ".toml", ".ini", ".cfg", ".log", ".sh", ".bat", ".ps1",
        ".html", ".htm", ".css", ".scss", ".less", ".svg", ".vue", ".svelte", ".astro",
    }

    out = []
    total_chars = 0
    with zipfile.ZipFile(io.BytesIO(raw)) as zf:
        infos = [i for i in zf.infolist() if not i.is_dir()]
        infos.sort(key=lambda i: (i.file_size, i.filename))
        picked = 0
        for info in infos:
            if picked >= max_entries or total_chars >= max_total_chars:
                break
            inner_ext = os.path.splitext(info.filename)[1].lower()
            if inner_ext not in text_like_ext and inner_ext not in {".pdf", ".docx", ".xlsx", ".pptx"}:
                continue
            if info.file_size > max_each_bytes:
                out.append(f"## {info.filename}\n[已跳过：文件过大 {info.file_size} bytes]")
                picked += 1
                continue
            try:
                blob = zf.read(info)
            except Exception:
                continue
            try:
                if inner_ext in text_like_ext:
                    inner_text = read_text_file(blob)
                elif inner_ext == ".pdf":
                    inner_text = read_pdf(blob)
                elif inner_ext == ".docx":
                    inner_text = read_docx(blob)
                elif inner_ext == ".xlsx":
                    inner_text = read_xlsx(blob)
                elif inner_ext == ".pptx":
                    inner_text = read_pptx(blob)
                else:
                    inner_text = ""
            except Exception as e:
                inner_text = f"[解析失败：{type(e).__name__}: {e}]"
            inner_text = truncate_text(inner_text or "", max_chars=max(400, min(4000, max_total_chars - total_chars)))
            if not inner_text:
                continue
            block = f"## {info.filename}\n{inner_text}"
            out.append(block)
            total_chars += len(block)
            picked += 1
    return "\n\n".join(out).strip()


OCR_ENABLE = app_getenv("OCR_ENABLE", "1").strip() != "0"
OCR_MAX_PAGES = max(1, int(app_getenv("OCR_MAX_PAGES", "10") or 10))
OCR_MAX_DOC_IMAGES = max(1, int(app_getenv("OCR_MAX_DOC_IMAGES", "10") or 10))
OCR_LANG = (app_getenv("OCR_LANG", "chi_sim+eng") or "chi_sim+eng").strip()
OCR_DPI = max(72, int(app_getenv("OCR_DPI", "300") or 300))


def _find_tesseract_cmd() -> str | None:
    env_path = (app_getenv("TESSERACT_CMD") or "").strip()
    if env_path and os.path.exists(env_path):
        return env_path
    for exe_name in ("tesseract.exe", "tesseract"):
        hit = shutil.which(exe_name)
        if hit:
            return hit
    candidates = []
    if os.name == "nt":
        candidates.extend([
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe"),
            os.path.expandvars(r"%USERPROFILE%\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"),
        ])
    for path in candidates:
        path = (path or "").strip()
        if path and os.path.exists(path):
            return path
    return None


def _find_poppler_path() -> str | None:
    env_path = (app_getenv("POPPLER_PATH") or "").strip()
    if env_path and os.path.isdir(env_path):
        return env_path
    for exe_name in ("pdftoppm.exe", "pdftoppm"):
        hit = shutil.which(exe_name)
        if hit:
            return os.path.dirname(hit)
    candidates = []
    if os.name == "nt":
        pf = os.environ.get("ProgramFiles", r"C:\Program Files")
        pf86 = os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")
        user = os.environ.get("USERPROFILE", "")
        candidates.extend([
            os.path.join(pf, "poppler", "Library", "bin"),
            os.path.join(pf86, "poppler", "Library", "bin"),
            os.path.join(user, "poppler", "Library", "bin"),
            r"D:\poppler\Library\bin",
            r"C:\poppler\Library\bin",
        ])
    for path in candidates:
        path = (path or "").strip()
        exe = "pdftoppm.exe" if os.name == "nt" else "pdftoppm"
        if path and os.path.isdir(path) and os.path.exists(os.path.join(path, exe)):
            return path
    return None


_OCR_EXTRA_VALID_SYMBOLS = '∨∧→¬↔↑↓≤≥≠≈±×÷∞∈∉⊆⊂⊇⊃∪∩∀∃⊥⊤∴∵⋯…△□◇○∥⊕⊗∑∏√≡∝∠∫∬∭∂∇⋮⋱⊢⊨↦⇔⇒⇐≌⊙⊖⊘⊚∦'


def _looks_symbol_dense_text(text: str) -> bool:
    s = str(text or '').strip()
    if not s:
        return False
    compact = re.sub(r'\s+', '', s)
    if not compact:
        return False
    symbol_hits = len(re.findall(rf'[{re.escape(_OCR_EXTRA_VALID_SYMBOLS)}]', s))
    if symbol_hits >= 2:
        return True
    if symbol_hits >= 1 and re.search(r'\b[pqra-z]\b', s, flags=re.I):
        return True
    return (float(symbol_hits) / float(max(1, len(compact)))) >= 0.03


def _ocr_text_score(text: str) -> float:
    s = str(text or '').strip()
    if not s:
        return 0.0
    compact = re.sub(r'\s+', '', s)
    if not compact:
        return 0.0
    allowed_punct = '，。！？、；：,.!?%()（）【】《》“”‘’:/_+\\-=#@&*[]{}<>|^~;\'"'
    allowed_chars = re.escape(allowed_punct + _OCR_EXTRA_VALID_SYMBOLS)
    valid_chars = re.findall(rf'[\u4e00-\u9fffA-Za-z0-9{allowed_chars}]', s)
    weird_chars = re.findall(rf'[^\s\u4e00-\u9fffA-Za-z0-9{allowed_chars}]', s)
    symbol_hits = len(re.findall(rf'[{re.escape(_OCR_EXTRA_VALID_SYMBOLS)}]', s))
    repeated_noise = len(re.findall(r'([A-Za-z0-9])\1{4,}', s))
    line_count = len([ln for ln in re.split(r'\r?\n+', s) if ln.strip()])
    score = float(len(compact))
    score += len(valid_chars) * 0.35
    score += symbol_hits * 0.75
    score += min(line_count, 12) * 1.5
    if _looks_symbol_dense_text(s):
        score += 12.0
    score -= len(weird_chars) * 1.1
    score -= repeated_noise * 8.0
    return max(0.0, score)



def _ocr_resize_for_text_detail(img):
    try:
        w, h = img.size
    except Exception:
        return img
    longest = max(int(w or 0), int(h or 0), 1)
    shortest = max(min(int(w or 0), int(h or 0)), 1)
    target_long = 2600 if shortest < 1100 else 2200
    if longest >= target_long:
        return img
    scale = min(2.8, max(1.0, float(target_long) / float(longest)))
    if scale <= 1.01:
        return img
    try:
        from PIL import Image  # type: ignore
        resample = getattr(Image, 'Resampling', Image).LANCZOS
    except Exception:
        return img
    try:
        return img.resize((max(1, int(round(w * scale))), max(1, int(round(h * scale)))), resample)
    except Exception:
        return img


def _ocr_estimate_background_gray(gray) -> int:
    try:
        from PIL import ImageStat  # type: ignore
    except Exception:
        return 255
    try:
        w, h = gray.size
        band = max(2, int(round(min(w, h) * 0.03)))
        strips = [
            gray.crop((0, 0, w, band)),
            gray.crop((0, max(0, h - band), w, h)),
            gray.crop((0, 0, band, h)),
            gray.crop((max(0, w - band), 0, w, h)),
        ]
        values = []
        for part in strips:
            try:
                values.append(float((ImageStat.Stat(part).mean or [255])[0]))
            except Exception:
                continue
        if not values:
            return 255
        values.sort()
        return max(0, min(255, int(round(values[len(values) // 2]))))
    except Exception:
        return 255


def _ocr_trim_background_edges(img, *, tolerance: int = 18, pad_ratio: float = 0.012):
    try:
        from PIL import ImageOps, ImageChops  # type: ignore
    except Exception:
        return img
    try:
        gray = ImageOps.grayscale(img)
        bg = _ocr_estimate_background_gray(gray)
        diff = ImageChops.difference(gray, ImageChops.constant(gray, int(bg)))
        mask = diff.point(lambda p: 255 if p >= max(6, int(tolerance or 0)) else 0, mode='L')
        bbox = mask.getbbox()
        if not bbox:
            return img
        w, h = gray.size
        left, top, right, bottom = bbox
        pad = max(4, int(round(min(w, h) * float(pad_ratio or 0.0))))
        left = max(0, left - pad)
        top = max(0, top - pad)
        right = min(w, right + pad)
        bottom = min(h, bottom + pad)
        if right - left < max(180, int(w * 0.38)) or bottom - top < max(140, int(h * 0.30)):
            return img
        if (right - left) >= int(w * 0.985) and (bottom - top) >= int(h * 0.985):
            return img
        return img.crop((left, top, right, bottom))
    except Exception:
        return img


def _ocr_render_text_ready_image(img, *, strong: bool = False):
    try:
        from PIL import ImageOps, ImageEnhance, ImageFilter  # type: ignore
    except Exception:
        return img
    try:
        gray = ImageOps.grayscale(img)
        gray = _ocr_resize_for_text_detail(gray)
        gray = ImageOps.autocontrast(gray, cutoff=1 if strong else 0)
        gray = ImageEnhance.Contrast(gray).enhance(2.05 if strong else 1.7)
        try:
            gray = gray.filter(ImageFilter.UnsharpMask(radius=2.0 if strong else 1.4, percent=235 if strong else 185, threshold=2))
        except Exception:
            pass
        gray = ImageEnhance.Sharpness(gray).enhance(2.35 if strong else 2.0)
        if strong:
            try:
                gray = gray.filter(ImageFilter.MedianFilter(size=3))
            except Exception:
                pass
        return gray
    except Exception:
        return img


def _ocr_build_adaptive_bw(img):
    try:
        from PIL import ImageOps, ImageFilter, ImageChops  # type: ignore
    except Exception:
        return None
    try:
        gray = ImageOps.grayscale(img)
        gray = _ocr_resize_for_text_detail(gray)
        gray = ImageOps.autocontrast(gray, cutoff=1)
        blurred = gray.filter(ImageFilter.BoxBlur(radius=6))
        local = ImageChops.add(gray, ImageChops.invert(blurred), scale=1.0, offset=118)
        bw = local.point(lambda p: 255 if p >= 132 else 0, mode='1').convert('L')
        return bw
    except Exception:
        return None


def _ocr_focus_region_score(img) -> float:
    try:
        from PIL import ImageOps, ImageStat  # type: ignore
    except Exception:
        return 0.0
    try:
        gray = ImageOps.grayscale(img)
        hist = gray.histogram() or []
        total = float(sum(hist) or 1.0)
        dark_ratio = float(sum(hist[:208])) / total
        mid_ratio = float(sum(hist[:244])) / total
        stat = ImageStat.Stat(gray)
        stddev = float((stat.stddev or [0.0])[0] or 0.0)
        w, h = gray.size
        area_bonus = min((float(w * h) / 240000.0), 12.0)
        return stddev * 2.2 + dark_ratio * 140.0 + mid_ratio * 28.0 + area_bonus
    except Exception:
        return 0.0


def _ocr_generate_focus_crops(img) -> list:
    try:
        from PIL import ImageOps  # type: ignore
    except Exception:
        return []
    out = []
    seen = set()

    def push(one):
        try:
            key = (tuple(getattr(one, 'size', (0, 0)) or (0, 0)), hash(one.tobytes()[:2048]))
        except Exception:
            key = (tuple(getattr(one, 'size', (0, 0)) or (0, 0)), id(one))
        if key in seen:
            return
        seen.add(key)
        out.append(one)

    try:
        base = img if getattr(img, 'mode', '') in ('RGB', 'L') else img.convert('RGB')
    except Exception:
        base = img
    trimmed = _ocr_trim_background_edges(base, tolerance=18)
    if getattr(trimmed, 'size', None) and getattr(trimmed, 'size', None) != getattr(base, 'size', None):
        push(trimmed)
    else:
        trimmed = base
    try:
        gray = ImageOps.grayscale(trimmed)
        bg = _ocr_estimate_background_gray(gray)
    except Exception:
        bg = 255
    try:
        w, h = trimmed.size
    except Exception:
        return out[:2]
    candidates = []
    def add_box(lf, tp, rf, bt):
        left = max(0, min(w - 1, int(round(w * float(lf)))))
        top = max(0, min(h - 1, int(round(h * float(tp)))))
        right = max(left + 1, min(w, int(round(w * float(rf)))))
        bottom = max(top + 1, min(h, int(round(h * float(bt)))))
        if right - left < max(220, int(w * 0.36)) or bottom - top < max(180, int(h * 0.26)):
            return
        try:
            crop = trimmed.crop((left, top, right, bottom))
        except Exception:
            return
        candidates.append(crop)

    add_box(0.03, 0.03, 0.97, 0.97)
    add_box(0.06, 0.06, 0.94, 0.94)
    if bg >= 238:
        if w >= int(h * 1.18) and h >= 620:
            add_box(0.14, 0.42, 0.92, 0.98)
            add_box(0.16, 0.50, 0.90, 0.98)
        if h >= int(w * 1.12) and w >= 520:
            add_box(0.08, 0.22, 0.92, 0.96)
            add_box(0.10, 0.34, 0.90, 0.98)
    scored = []
    for crop in candidates:
        score = _ocr_focus_region_score(crop)
        scored.append((score, crop))
    scored.sort(key=lambda it: it[0], reverse=True)
    for _score, crop in scored[:2]:
        push(crop)
    return out[:2]


def _ocr_prepare_variants(img) -> list:
    try:
        from PIL import ImageOps  # type: ignore
    except Exception:
        return [img]

    variants = []
    seen = set()

    def push(one):
        try:
            key = (getattr(one, 'mode', ''), tuple(getattr(one, 'size', (0, 0)) or (0, 0)), hash(one.tobytes()[:2048]))
        except Exception:
            key = (getattr(one, 'mode', ''), tuple(getattr(one, 'size', (0, 0)) or (0, 0)), id(one))
        if key in seen:
            return
        seen.add(key)
        variants.append(one)

    base = img
    if getattr(base, 'mode', '') not in ('RGB', 'L'):
        try:
            base = base.convert('RGB')
        except Exception:
            pass
    push(base)

    try:
        trimmed = _ocr_trim_background_edges(base, tolerance=18)
        push(trimmed)
        enhanced = _ocr_render_text_ready_image(trimmed, strong=False)
        push(enhanced)
        strong = _ocr_render_text_ready_image(trimmed, strong=True)
        push(strong)
        bw = _ocr_build_adaptive_bw(strong)
        if bw is not None:
            push(bw)
        for crop in _ocr_generate_focus_crops(base):
            focus_ready = _ocr_render_text_ready_image(crop, strong=True)
            push(focus_ready)
            focus_bw = _ocr_build_adaptive_bw(crop)
            if focus_bw is not None:
                push(focus_bw)
    except Exception:
        pass

    return variants[:6]


def _ocr_image_pil(img) -> str:
    if not OCR_ENABLE:
        return ""
    try:
        import pytesseract  # type: ignore
    except Exception:
        return ""
    try:
        tess_cmd = _find_tesseract_cmd()
        if tess_cmd:
            pytesseract.pytesseract.tesseract_cmd = tess_cmd
    except Exception:
        pass

    configs = ['--oem 3 --psm 6', '--oem 3 --psm 11']
    best_text = ''
    best_score = -1.0

    for idx, variant in enumerate(_ocr_prepare_variants(img)):
        cfgs = list(configs)
        if idx >= 1 and best_score >= 120:
            cfgs = cfgs[:1]
        for cfg in cfgs:
            text = ''
            try:
                text = (pytesseract.image_to_string(variant, lang=OCR_LANG, config=cfg) or '').strip()
            except Exception:
                try:
                    text = (pytesseract.image_to_string(variant, config=cfg) or '').strip()
                except Exception:
                    text = ''
            score = _ocr_text_score(text)
            if score > best_score:
                best_score = score
                best_text = text
        if best_score >= 220:
            break

    return best_text.strip()


def _ocr_image_bytes(raw: bytes) -> str:
    if not raw or not OCR_ENABLE:
        return ""
    try:
        from PIL import Image, ImageOps  # type: ignore
    except Exception:
        return ""
    try:
        img = Image.open(io.BytesIO(raw))
        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass
        if img.mode not in ("L", "RGB"):
            try:
                img = img.convert("RGB")
            except Exception:
                pass
        return _ocr_image_pil(img)
    except Exception:
        return ""


def _ocr_pdf(raw: bytes, max_pages: int | None = None) -> str:
    if not raw or not OCR_ENABLE:
        return ""
    try:
        from pdf2image import convert_from_bytes  # type: ignore
    except Exception:
        return ""
    pages = max(1, int(max_pages or OCR_MAX_PAGES))
    kwargs = {"dpi": OCR_DPI, "first_page": 1, "last_page": pages}
    poppler_path = _find_poppler_path()
    if poppler_path:
        kwargs["poppler_path"] = poppler_path
    try:
        images = convert_from_bytes(raw, **kwargs)
    except Exception:
        return ""
    out = []
    for img in images[:pages]:
        txt = _ocr_image_pil(img)
        if txt:
            out.append(txt)
    return "\n".join(out).strip()


def _iter_docx_image_blobs(raw: bytes, max_images: int | None = None):
    limit = max(1, int(max_images or OCR_MAX_DOC_IMAGES))
    try:
        import zipfile
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            names = [n for n in zf.namelist() if n.startswith("word/media/")]
            names.sort()
            count = 0
            for name in names:
                if count >= limit:
                    break
                try:
                    blob = zf.read(name)
                except Exception:
                    continue
                if blob:
                    count += 1
                    yield blob
    except Exception:
        return


def _ocr_docx_images(raw: bytes, max_images: int | None = None) -> str:
    out = []
    for blob in _iter_docx_image_blobs(raw, max_images=max_images):
        txt = _ocr_image_bytes(blob)
        if txt:
            out.append(txt)
    return "\n".join(out).strip()


def read_pdf(raw: bytes) -> str:
    """Extract PDF text layer only; visual/scanned pages use sandbox_analyze_file_images."""
    if not _looks_like_pdf_bytes(raw):
        raise ValueError("body_not_pdf")
    max_pdf_bytes = int(app_getenv("PDF_PARSE_MAX_BYTES", "20000000") or 20000000)
    if len(raw) > max_pdf_bytes:
        raw = raw[:max_pdf_bytes]
    try:
        reader = PdfReader(io.BytesIO(raw))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
        return "\n\n".join(parts).strip()
    except Exception:
        pass
    return ""


def read_docx(raw: bytes) -> str:
    try:
        doc = Document(io.BytesIO(raw))
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        return "\n".join(paras).strip()
    except Exception:
        pass
    return ""


def _sandbox_docx_text_diagnostics(raw: bytes, text: str = '', filename: str = '') -> dict:
    """Inspect DOCX OOXML signals that are lost when only paragraph text is read."""
    info: dict = {
        'document_type': 'docx',
        'filename': os.path.basename(str(filename or '')),
        'zip_entries': 0,
        'paragraph_count': 0,
        'nonempty_paragraph_count': 0,
        'table_count': 0,
        'office_math_count': 0,
        'media_count': 0,
        'media_by_ext': {},
        'paragraph_count': 0,
        'nonempty_paragraph_count': 0,
        'table_count': 0,
        'office_math_count': 0,
        'drawing_count': 0,
        'object_count': 0,
        'shape_count': 0,
        'ole_count': 0,
        'embedding_count': 0,
        'table_summaries': [],
        'caption_candidates': [],
        'text_gap_samples': [],
        'quality_signals': [],
        'requires_visual_review': False,
    }
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        ns = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
            'm': 'http://schemas.openxmlformats.org/officeDocument/2006/math',
            'v': 'urn:schemas-microsoft-com:vml',
            'o': 'urn:schemas-microsoft-com:office:office',
        }

        def norm(value) -> str:
            return re.sub(r'\s+', ' ', str(value or '').replace('\r', ' ')).strip()

        def compact(value) -> str:
            return re.sub(r'\s+', '', str(value or ''))

        def elem_text(elem) -> str:
            return norm(''.join((t.text or '') for t in elem.findall('.//w:t', ns)))

        with zipfile.ZipFile(io.BytesIO(raw), 'r') as zf:
            names = zf.namelist()
            info['zip_entries'] = len(names)
            media = [n for n in names if n.lower().startswith('word/media/') and not n.endswith('/')]
            info['media_count'] = len(media)
            by_ext: dict[str, int] = {}
            for name in media:
                ext = os.path.splitext(name)[1].lower() or '<none>'
                by_ext[ext] = by_ext.get(ext, 0) + 1
            info['media_by_ext'] = by_ext
            info['embedding_count'] = len([n for n in names if n.lower().startswith('word/embeddings/') and not n.endswith('/')])
            if 'word/document.xml' not in names:
                info['quality_signals'].append('docx_missing_word_document_xml')
                info['requires_visual_review'] = True
                return info
            root = ET.fromstring(zf.read('word/document.xml'))
            paras = root.findall('.//w:p', ns)
            info['paragraph_count'] = len(paras)
            para_texts = [elem_text(p) for p in paras]
            info['nonempty_paragraph_count'] = len([x for x in para_texts if x])
            info['table_count'] = len(root.findall('.//w:tbl', ns))
            info['office_math_count'] = len(root.findall('.//m:oMath', ns)) + len(root.findall('.//m:oMathPara', ns))
            info['drawing_count'] = len(root.findall('.//w:drawing', ns))
            info['object_count'] = len(root.findall('.//w:object', ns))
            info['shape_count'] = len(root.findall('.//v:shape', ns))
            info['ole_count'] = len(root.findall('.//o:OLEObject', ns))

            body = root.find('.//w:body', ns)
            recent_paras: list[str] = []
            table_summaries: list[dict] = []
            captions: list[dict] = []
            if body is not None:
                for child in list(body):
                    tag = str(child.tag or '').rsplit('}', 1)[-1]
                    if tag == 'p':
                        t = elem_text(child)
                        if t:
                            recent_paras.append(t)
                            recent_paras = recent_paras[-5:]
                            c = compact(t)
                            if (re.search(r'(\u56fe|\u8868)\s*[0-9\u4e00-\u9fff]+', t) or re.search(r'(Figure|Fig\.?|Table)\s*[0-9]+', t, re.I)) and len(captions) < 24:
                                captions.append({'text': t[:260], 'near_table': False})
                    elif tag == 'tbl':
                        rows: list[list[str]] = []
                        for tr in child.findall('.//w:tr', ns):
                            cells = [elem_text(tc) for tc in tr.findall('./w:tc', ns)]
                            if cells:
                                rows.append(cells)
                        col_count = max([len(r) for r in rows] or [0])
                        empty_by_col = []
                        for col_idx in range(min(col_count, 8)):
                            vals = [(row[col_idx] if col_idx < len(row) else '') for row in rows]
                            empty_by_col.append(len([v for v in vals if not norm(v)]))
                        first_col_empty_ratio = (float(empty_by_col[0]) / float(max(1, len(rows)))) if empty_by_col else 0.0
                        title_context = ' / '.join(recent_paras[-3:])
                        summary = {
                            'index': len(table_summaries) + 1,
                            'rows': len(rows),
                            'cols': col_count,
                            'title_context': title_context[:300],
                            'first_col_empty_ratio': round(first_col_empty_ratio, 3),
                            'empty_cells_first_columns': empty_by_col,
                            'sample_rows': [[cell[:80] for cell in row[:4]] for row in rows[:4]],
                        }
                        if first_col_empty_ratio >= 0.65 and len(rows) >= 3:
                            summary['issue'] = 'leading_column_mostly_empty'
                            info['quality_signals'].append('table_leading_column_mostly_empty')
                        table_summaries.append(summary)
                        if len(table_summaries) >= 8:
                            continue
            info['table_summaries'] = table_summaries[:8]
            info['caption_candidates'] = captions[:24]

            lines = [norm(x) for x in str(text or '').splitlines() if norm(x)]
            gap_patterns = [
                ('bare_current_year_explanation', r'表示现在[，,]?\s*表示第?1?年末到第?5?年末'),
                ('dangling_step', r'^第一步[，,。；;:]?$'),
                ('dangling_then', r'^此时[，,。；;:]?$'),
                ('missing_theory_cost_formula', r'若理论花费.*实际需求等于理论需求[：:；;]?$'),
                ('missing_subscripts_year_area_class', r'第年末第个小区第类老人'),
                ('dangling_formula_colon', r'[：:；;]\s*$'),
            ]
            gap_samples: list[dict] = []
            for idx, line in enumerate(lines, 1):
                for code, pat in gap_patterns:
                    if re.search(pat, line):
                        gap_samples.append({'line': idx, 'type': code, 'text': line[:220]})
                        break
                if len(gap_samples) >= 16:
                    break
            info['text_gap_samples'] = gap_samples
            if gap_samples:
                info['quality_signals'].append('text_layer_has_formula_or_symbol_gaps')
            vector_count = int(by_ext.get('.wmf') or 0) + int(by_ext.get('.emf') or 0)
            if vector_count >= 8 and int(info.get('office_math_count') or 0) == 0:
                info['quality_signals'].append('many_vector_media_but_no_office_math')
            if int(info.get('media_count') or 0) > 0 and (gap_samples or vector_count >= 8):
                info['requires_visual_review'] = True
            if int(info.get('table_count') or 0) and any((x.get('issue') == 'leading_column_mostly_empty') for x in table_summaries):
                info['requires_visual_review'] = True
    except Exception as exc:
        info['inspect_error'] = f'{type(exc).__name__}: {exc}'
    return info


def _sandbox_document_text_diagnostics(raw: bytes, filename: str = '', text: str = '') -> dict:
    ext = os.path.splitext(str(filename or ''))[1].lower()
    if ext == '.docx':
        return _sandbox_docx_text_diagnostics(raw, text=text, filename=filename)
    return {}


def _sandbox_document_diagnostic_summary(diagnostics: dict | None = None) -> dict:
    diag = diagnostics if isinstance(diagnostics, dict) else {}
    if not diag:
        return {}
    tables = []
    for item in (diag.get('table_summaries') or [])[:4]:
        if not isinstance(item, dict):
            continue
        tables.append({
            'index': item.get('index'),
            'rows': item.get('rows'),
            'cols': item.get('cols'),
            'title_context': str(item.get('title_context') or '')[:220],
            'first_col_empty_ratio': item.get('first_col_empty_ratio'),
            'issue': str(item.get('issue') or '')[:120],
            'sample_rows': item.get('sample_rows')[:3] if isinstance(item.get('sample_rows'), list) else [],
        })
    return {
        'document_type': str(diag.get('document_type') or '')[:40],
        'paragraph_count': int(diag.get('paragraph_count') or 0),
        'nonempty_paragraph_count': int(diag.get('nonempty_paragraph_count') or 0),
        'table_count': int(diag.get('table_count') or 0),
        'office_math_count': int(diag.get('office_math_count') or 0),
        'media_count': int(diag.get('media_count') or 0),
        'media_by_ext': diag.get('media_by_ext') if isinstance(diag.get('media_by_ext'), dict) else {},
        'drawing_count': int(diag.get('drawing_count') or 0),
        'object_count': int(diag.get('object_count') or 0),
        'shape_count': int(diag.get('shape_count') or 0),
        'ole_count': int(diag.get('ole_count') or 0),
        'embedding_count': int(diag.get('embedding_count') or 0),
        'requires_visual_review': bool(diag.get('requires_visual_review')),
        'quality_signals': [str(x or '')[:120] for x in (diag.get('quality_signals') or [])[:12]],
        'text_gap_samples': [dict(x) for x in (diag.get('text_gap_samples') or [])[:8] if isinstance(x, dict)],
        'table_summaries': tables,
        'caption_candidates': [dict(x) for x in (diag.get('caption_candidates') or [])[:8] if isinstance(x, dict)],
        'inspect_error': str(diag.get('inspect_error') or '')[:260],
    }


def read_xlsx(raw: bytes, max_sheets=3, max_rows=80, max_cols=25) -> str:
    """Read spreadsheet files as structured text first.

    This is the default evidence lane for XLSX. Rendered page images are only a
    secondary lane for explicit chart/layout/format/merged-cell/page questions.
    """
    wb_values = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=False)
    try:
        wb_formula = openpyxl.load_workbook(io.BytesIO(raw), data_only=False, read_only=False)
    except Exception:
        wb_formula = None
    out = [
        "# Workbook summary",
        f"sheets={len(wb_values.sheetnames)} names={', '.join(wb_values.sheetnames[:max_sheets])}",
    ]
    for sheet_name in wb_values.sheetnames[:max_sheets]:
        ws = wb_values[sheet_name]
        wf = wb_formula[sheet_name] if wb_formula is not None and sheet_name in wb_formula.sheetnames else None
        max_row = int(ws.max_row or 0)
        max_col = int(ws.max_column or 0)
        try:
            merged_count = len(list(ws.merged_cells.ranges))
        except Exception:
            merged_count = 0
        chart_count = len(getattr(ws, '_charts', []) or [])
        image_count = len(getattr(ws, '_images', []) or [])
        out.append("")
        out.append(f"# Sheet: {sheet_name}")
        out.append(f"[sheet_meta rows={max_row} cols={max_col} merged_cells={merged_count} charts={chart_count} images={image_count}]")
        if max_row > max_rows or max_col > max_cols:
            out.append(f"[preview limited to rows<= {max_rows}, cols<= {max_cols}]")
        for r in range(1, min(max_row, max_rows) + 1):
            vals = []
            for c in range(1, min(max_col, max_cols) + 1):
                v = ws.cell(row=r, column=c).value
                cell_text = "" if v is None else str(v)
                if wf is not None:
                    try:
                        fv = wf.cell(row=r, column=c).value
                        if isinstance(fv, str) and fv.startswith('=') and fv != cell_text:
                            cell_text = (cell_text + f" [formula: {fv}]").strip()
                    except Exception:
                        pass
                vals.append(cell_text)
            if max_col > max_cols:
                vals.append("[…]")
            out.append("\t".join(vals).rstrip())
        if max_row > max_rows:
            out.append("[已截断：行数过多]")
    for wb in (wb_values, wb_formula):
        try:
            if wb is not None:
                wb.close()
        except Exception:
            pass
    return "\n".join(out).strip()


def read_xls(raw: bytes) -> str:
    # 需要 pandas + xlrd
    import pandas as pd  # type: ignore

    with tempfile.NamedTemporaryFile(suffix=".xls", delete=False) as tf:
        tf.write(raw)
        temp_xls = tf.name
    try:
        xls = pd.ExcelFile(temp_xls, engine="xlrd")
        out = []
        for sheet_name in xls.sheet_names[:3]:
            df = pd.read_excel(xls, sheet_name=sheet_name, header=None, nrows=80)
            if df.shape[1] > 25:
                df = df.iloc[:, :25]
            out.append(f"# Sheet: {sheet_name}")
            out.append(df.fillna("").astype(str).to_csv(sep="\t", index=False, header=False).strip())
            out.append("")
        return "".join(out).strip()
    finally:
        try:
            os.unlink(temp_xls)
        except Exception:
            pass


def _find_soffice() -> str | None:
    for name in ("soffice", "soffice.exe"):
        p = shutil.which(name)
        if p:
            return p
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    return None


def _soffice_convert(input_path: str, out_dir: str, to_ext: str) -> str:
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError("未找到 LibreOffice/soffice")

    cmd = [
        soffice,
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to", to_ext,
        "--outdir", out_dir,
        input_path,
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"soffice 转换失败：{proc.stderr.strip() or proc.stdout.strip() or 'unknown error'}")

    base = os.path.splitext(os.path.basename(input_path))[0]
    out_path = os.path.join(out_dir, f"{base}.{to_ext}")
    if os.path.exists(out_path):
        return out_path

    for fn in os.listdir(out_dir):
        if fn.lower().endswith("." + to_ext.lower()):
            return os.path.join(out_dir, fn)

    raise RuntimeError("soffice 转换后未找到输出文件")


def _convert_doc_to_docx_via_word(doc_path: str, out_dir: str) -> str:
    """
    Windows 下用 Microsoft Word COM 把 .doc 转成 .docx
    需要安装 Microsoft Word + pywin32
    修复：多线程/Flask 环境下必须 CoInitialize
    """
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except Exception as e:
        raise RuntimeError("需要 pywin32（py -m pip install pywin32）") from e

    word = None
    doc = None
    pythoncom.CoInitialize()
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0

        doc = word.Documents.Open(doc_path, ReadOnly=True)
        base = os.path.splitext(os.path.basename(doc_path))[0]
        out_path = os.path.join(out_dir, f"{base}.docx")

        # 16 => wdFormatXMLDocument (.docx)
        doc.SaveAs2(out_path, FileFormat=16)
        doc.Close(False)
        doc = None

        word.Quit()
        word = None

        if not os.path.exists(out_path):
            raise RuntimeError("Word 转换后未找到输出 docx")
        return out_path
    except Exception as e:
        raise RuntimeError(f"Word 转换失败: {type(e).__name__}: {e}") from e
    finally:
        try:
            if doc is not None:
                doc.Close(False)
        except Exception:
            pass
        try:
            if word is not None:
                word.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def read_doc(raw: bytes) -> str:
    """
    .doc -> 先尝试 Word COM 转 docx，再用 python-docx 解析
    若 Word 失败：只有在 soffice 存在时才 fallback，否则把 Word 失败原因直接抛出
    """
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, "upload.doc")
        with open(in_path, "wb") as f:
            f.write(raw)

        try:
            out_docx = _convert_doc_to_docx_via_word(in_path, td)
        except Exception as e_word:
            soffice = _find_soffice()
            if soffice:
                out_docx = _soffice_convert(in_path, td, "docx")
            else:
                raise RuntimeError(f"Word 转换失败：{type(e_word).__name__}: {e_word}；并且未找到 LibreOffice/soffice") from e_word

        with open(out_docx, "rb") as f:
            return read_docx(f.read())
