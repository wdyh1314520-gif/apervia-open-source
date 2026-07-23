import argparse
import importlib
import json
import os
import subprocess
import sys
from pathlib import Path


MODULES = [
    "bs4",
    "cairosvg",
    "cv2",
    "docx",
    "fitz",
    "lxml",
    "matplotlib",
    "numpy",
    "openpyxl",
    "pandas",
    "pdfplumber",
    "PIL",
    "playwright",
    "pptx",
    "pypdf",
    "pytesseract",
    "reportlab",
    "skimage",
    "weasyprint",
    "xlsxwriter",
]


def run(cmd):
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return {
        "cmd": cmd,
        "exit_code": proc.returncode,
        "stdout": (proc.stdout or "").strip().splitlines()[:3],
        "stderr": (proc.stderr or "").strip().splitlines()[:3],
    }


def import_check():
    ok = []
    missing = []
    for name in MODULES:
        try:
            importlib.import_module(name)
            ok.append(name)
        except Exception as exc:
            missing.append({"module": name, "error": f"{type(exc).__name__}: {exc}"})
    return ok, missing


def file_check(root):
    root.mkdir(parents=True, exist_ok=True)

    from docx import Document
    from openpyxl import Workbook
    from PIL import Image, ImageDraw
    from pptx import Presentation
    from reportlab.pdfgen import canvas

    doc = Document()
    doc.add_paragraph("app3 sandbox docx probe")
    doc.save(root / "probe.docx")

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "app3"
    ws["B1"] = 312
    wb.save(root / "probe.xlsx")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "app3 sandbox pptx probe"
    prs.save(root / "probe.pptx")

    pdf_path = root / "probe.pdf"
    c = canvas.Canvas(str(pdf_path))
    c.drawString(72, 720, "app3 sandbox pdf probe")
    c.save()

    img = Image.new("RGB", (360, 120), "white")
    draw = ImageDraw.Draw(img)
    draw.text((20, 40), "app3 OCR probe", fill="black")
    img.save(root / "probe.png")

    return sorted(p.name for p in root.glob("probe.*"))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--quick", action="store_true")
    args = parser.parse_args()

    ok, missing = import_check()
    out = {
        "python": sys.version.split()[0],
        "import_ok_count": len(ok),
        "missing": missing,
        "tools": {
            "tesseract": run(["tesseract", "--version"]),
            "pdftoppm": run(["pdftoppm", "-v"]),
            "libreoffice": run(["libreoffice", "--version"]),
            "node": run(["node", "--version"]),
        },
    }
    if not args.quick:
        out["files"] = file_check(Path(os.environ.get("APP3_PROBE_DIR", "/mnt/data/probe")))
    print(json.dumps(out, ensure_ascii=False, indent=2))
    if missing:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
